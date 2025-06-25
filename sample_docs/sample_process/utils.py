from typing import List
import PyPDF2
import re
from langdetect import detect
from docx import Document
from pptx import Presentation

from mongodb.ops.object_op import get_objects_by_conditions


# 解析 PDF 文件为文本
def extract_text_from_pdf(file_path: str) -> List[str]:
    """读取 PDF 并返回按页分割的文本列表"""
    texts = []
    try:
        with open(file_path, "rb") as file:
            reader = PyPDF2.PdfReader(file)
            for page in reader.pages:
                texts.append(page.extract_text() or "")  # 处理可能为空的页面
    except Exception as e:
        print(f"解析 PDF 失败: {file_path}, 错误: {e}")
    return texts


def remove_vietnamese_text(text: str) -> str:
    """移除越南语文本，并将多个连续换行符合并为一个"""
    # 1. 删除包含越南语字符的单词
    VIETNAMESE_CHAR_PATTERN = r"[ăâđêôơưĂÂĐÊÔƠƯ]"
    cleaned_text = re.sub(VIETNAMESE_CHAR_PATTERN, "", text)

    # 2. 按行处理语言检测，移除越南语句子
    sentences = cleaned_text.split("\n")
    filtered_sentences = []

    for sentence in sentences:
        try:
            if detect(sentence.strip()) != "vi":  # 过滤越南语
                filtered_sentences.append(sentence.strip())  # 移除首尾空格
        except:
            filtered_sentences.append(sentence.strip())  # 短句保留

    # 3. 重新拼接，并去掉多余的换行
    cleaned_text = "\n".join(filtered_sentences).strip()

    # 4. 将多个连续换行符合并为一个
    cleaned_text = re.sub(r'\n+', '\n', cleaned_text)

    return cleaned_text


def count_words(text: str) -> (int, List[str]):
    """改进 word_count 计算：中文按字符计数，英文按单词计数"""
    words = []
    text = text.replace('\n', " ")
    for sentence in text.split(" "):
        try:
            lang = detect(sentence)
            # print(sentence, lang)
            if lang == "zh-cn":  # 中文
                words.extend(list(sentence))  # 逐字符添加
            else:  # 其他语言（如英文）
                words.extend(sentence.split())  # 按单词分割
        except:
            words.extend(sentence.split())  # 保底策略
    # print(words)
    return len(words), words


def extract_text_from_docx(file_path: str) -> List[str]:
    """读取 Word 并返回按段落分割的文本列表"""
    texts = []
    try:
        document = Document(file_path)
        for paragraph in document.paragraphs:
            texts.append(paragraph.text)
    except Exception as e:
        print(f"解析 Word 失败: {file_path}, 错误: {e}")
    return texts


from typing import List
from pptx import Presentation


def extract_text_from_pptx(file_path: str) -> List[str]:
    """读取 PowerPoint 并返回按幻灯片分割的文本列表，包括备注"""
    texts = []
    try:
        presentation = Presentation(file_path)
        for slide in presentation.slides:
            slide_text = []

            # 提取幻灯片文本
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    slide_text.append(shape.text.strip())

            # 提取备注文本
            if slide.has_notes_slide:
                notes_slide = slide.notes_slide
                if notes_slide.notes_text_frame and notes_slide.notes_text_frame.text:
                    slide_text.append("备注: " + notes_slide.notes_text_frame.text.strip())

            texts.append("\n".join(slide_text))
    except Exception as e:
        print(f"解析 PowerPoint 失败: {file_path}, 错误: {e}")
    return texts


if __name__ == "__main__":
    # ss = extract_text_from_pdf("../人力资源相关课件及管理办法/01-立讯精密培训管理办法.pdf")
    # for p in texts:
    #     print(p)
    # print(sum([len(s) for s in ss]))
    # txt = "立讯公司员工离职需确认保密协议，COC 考试结果影响晋升、薪资及绩效。HR 负责培训、考试监督，未完成考试者受限，HR 中心执行整改及处罚。附"
    # print(count_words(txt))
    path = "/root/ns-rag/sample_docs/batch2/LXJS-QWI-GM-01__1_立讯技术公务接待管理办法__1_.pdf"
    # path = "/root/ns-rag/sample_docs/batch2/LXEP-026_1__防止歧视及骚扰作业程序.doc"
    print(extract_text_from_pdf(path))